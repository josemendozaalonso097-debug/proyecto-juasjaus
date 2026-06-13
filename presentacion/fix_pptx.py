from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

path_in  = 'attached_assets/Presentación_de_Checoso_20260613_124634_0000_1781376402468.pptx'
path_out = 'presentacion/CBTis258_Servicios_Financieros_Final.pptx'

prs = Presentation(path_in)
slides = list(prs.slides)

# ─── helpers ───────────────────────────────────────────────────────────────

def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

def replace_text(shape, new_text, font_size=None, bold=None, color=None, align=None):
    """Reemplaza TODO el texto de un textbox preservando la estructura de párrafos."""
    tf = shape.text_frame
    # Limpiar todos los párrafos existentes menos el primero
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]
        p._p.getparent().remove(p._p)
    lines = new_text.split('\n') if isinstance(new_text, str) else new_text
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
            # Limpiar runs existentes
            for run in p.runs:
                run.text = ''
            if p.runs:
                p.runs[0].text = line
                run = p.runs[0]
            else:
                run = p.add_run()
                run.text = line
        else:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
        if font_size:
            run.font.size = Pt(font_size)
        if bold is not None:
            run.font.bold = bold
        if color:
            run.font.color.rgb = RGBColor(*color)
        if align is not None:
            p.alignment = align

def add_box(slide, left, top, width, height, lines,
            font_size=14, bold=False, color=(255,255,255),
            align=PP_ALIGN.LEFT, bg_color=None, word_wrap=True):
    """Agrega un nuevo textbox al slide."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    for i, line in enumerate(lines if isinstance(lines, list) else [lines]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
        p.alignment = align
    return txBox

WHITE = (255,255,255)
RED   = (242, 13, 13)
DARK  = (30, 30, 30)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Portada
# ══════════════════════════════════════════════════════════════════════════════
s1 = slides[0]

# Materia: "IBDNRSi" → correcto
tb_materia = find_shape(s1, 'TextBox 33')
if tb_materia:
    replace_text(tb_materia, 'Desarrollo de Aplicaciones Web', font_size=18, color=WHITE)

# Slogan: "Un motivo de orgullo" → slogan real
tb_slogan = find_shape(s1, 'TextBox 31')
if tb_slogan:
    replace_text(tb_slogan, 'Digitalización de servicios financieros y administrativos', font_size=14, color=WHITE)

# Agregar badges de stack en la parte inferior izquierda
add_box(s1, 0.8, 9.5, 9.0, 0.6,
        'Python · FastAPI · React · SQLite · JWT · API REST · Google OAuth',
        font_size=12, bold=False, color=(200,200,200))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — El Problema + Nuestra Solución
# ══════════════════════════════════════════════════════════════════════════════
s2 = slides[1]

# Título columna derecha
add_box(s2, 12.4, 0.4, 6.5, 0.8,
        'NUESTRA SOLUCIÓN', font_size=22, bold=True, color=RED)

# Bloque: Innovación
add_box(s2, 12.4, 1.4, 6.5, 0.45,
        '💡 Innovación', font_size=15, bold=True, color=WHITE)
add_box(s2, 12.4, 1.9, 6.5, 0.8,
        'La única plataforma del plantel que unifica pagos, tienda, trámites y eventos en un solo lugar accesible desde cualquier dispositivo.',
        font_size=12, bold=False, color=(210,210,210))

# Bloque: Solución Digital
add_box(s2, 12.4, 3.0, 6.5, 0.45,
        '🖥️ Solución Digital', font_size=15, bold=True, color=WHITE)
add_box(s2, 12.4, 3.5, 6.5, 0.8,
        'Digitaliza el pago de colegiaturas, compra de uniformes y libros, solicitud de trámites y consulta de eventos — sin filas ni dinero físico.',
        font_size=12, bold=False, color=(210,210,210))

# Bloque: Usuario Meta
add_box(s2, 12.4, 4.6, 6.5, 0.45,
        '🎯 Usuario Meta', font_size=15, bold=True, color=WHITE)
add_box(s2, 12.4, 5.1, 6.5, 0.6,
        'Alumnos y padres de familia del CBTis 258 que necesitan acceder a los servicios escolares sin ir físicamente a ventanilla.',
        font_size=12, bold=False, color=(210,210,210))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Arquitectura  (fix "mongo" → "ORM")
# ══════════════════════════════════════════════════════════════════════════════
s3 = slides[2]
tb_arch = find_shape(s3, 'TextBox 16')
if tb_arch:
    replace_text(tb_arch,
        'Capa                            Tecnología\n'
        'Frontend                   React 19, Tailwind CSS, Vite\n'
        'Backend                   Python 3.12, FastAPI, Uvicorn\n'
        'Base de datos         SQLite + SQLAlchemy ORM\n'
        'Autenticación          JWT + Google OAuth 2.0\n'
        'Email                         Gmail SMTP (códigos OTP)',
        font_size=14, color=WHITE)

# Agregar fila de badges
add_box(s3, 1.5, 8.8, 17.0, 0.55,
        '🔐 JWT Auth   ·   ⚡ Axios Async/Await   ·   🌐 API REST (JSON)   ·   📦 SQLAlchemy ORM   ·   📧 Gmail SMTP',
        font_size=13, bold=False, color=(200,200,200), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Base de Datos  (MongoDB → SQLite, texto → esquemas)
# ══════════════════════════════════════════════════════════════════════════════
s4 = slides[3]

# Cambiar subtítulo "MongoDB" → "SQLite"
tb_mongo = find_shape(s4, 'TextBox 20')
if tb_mongo:
    replace_text(tb_mongo, 'SQLite', font_size=28, bold=True, color=RED)

# Reemplazar texto descriptivo con esquemas de campos
tb_desc = find_shape(s4, 'TextBox 21')
if tb_desc:
    replace_text(tb_desc,
        '{ users }                              { productos }                    { eventos }\n'
        '_id: Integer                           _id: Integer                      _id: Integer\n'
        'nombre: String                       nombre: String                  titulo: String\n'
        'email: String (único)             marca: String                    fecha: String\n'
        'password: String (hash)       precio: Float                      descripcion: Text\n'
        'rol: String                              categoria: String               createdBy: Integer\n'
        'semestre: String                     imagen: String\n'
        '\n'
        '{ otp_codes }                         { compras }                      { password_resets }\n'
        '_id: Integer                           _id: Integer                      _id: Integer\n'
        'email: String                          user_id: Integer               user_id: Integer\n'
        'codigo: String                        total: Float                       token: String (único)\n'
        'expira: DateTime                    fecha: DateTime               expira: DateTime',
        font_size=11, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — API REST Endpoints  (párrafos → tabla)
# ══════════════════════════════════════════════════════════════════════════════
s5 = slides[4]
tb_api = find_shape(s5, 'TextBox 27')
if tb_api:
    replace_text(tb_api,
        'Método   Endpoint                                Descripción                          Auth\n'
        '──────────────────────────────────────────────────────────────────────────────────────\n'
        'POST      /api/auth/register-send-otp     Envía código OTP al correo          —\n'
        'POST      /api/auth/register-verify-otp   Verifica OTP y crea cuenta           —\n'
        'POST      /api/auth/login                         Login con email y contraseña     —\n'
        'POST      /api/auth/login/google             Login con cuenta de Google          —\n'
        'POST      /api/auth/forgot-password        Envía email de recuperación          —\n'
        'POST      /api/auth/reset-password          Cambia contraseña con token         —\n'
        'GET         /api/auth/check-session            Valida token JWT activo               🔐\n'
        'GET         /api/tienda/productos               Lista productos del catálogo         —\n'
        'POST      /api/tienda/productos               Crea nuevo producto                     🔐 Admin\n'
        'PUT          /api/tienda/productos/:id        Edita un producto                          🔐 Admin\n'
        'DELETE    /api/tienda/productos/:id        Elimina un producto                       🔐 Admin\n'
        'GET         /api/eventos                              Lista eventos escolares               🔐\n'
        'POST      /api/eventos                              Crea un evento                              🔐 Admin\n'
        'PUT          /api/eventos/:id                       Edita un evento                              🔐 Admin\n'
        'DELETE    /api/eventos/:id                       Elimina un evento                           🔐 Admin',
        font_size=10, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Funcionalidades Clave  (agregar los 3 faltantes)
# ══════════════════════════════════════════════════════════════════════════════
s8 = slides[7]
tb_func = find_shape(s8, 'TextBox 28')
if tb_func:
    replace_text(tb_func,
        'Para alumnos:\n'
        '🔐 Registro con verificación OTP de 6 dígitos por correo\n'
        '🛒 Tienda escolar: uniformes, libros, credencial, trámites\n'
        '💳 Pago simulado: tarjeta, OXXO (código de barras), transferencia\n'
        '📄 Generación de recibos en PDF automáticamente\n'
        '📅 Consultar eventos y avisos escolares\n'
        '👤 Perfil personal editable con foto\n'
        '\n'
        'Para administradores:\n'
        '📦 Gestión completa del catálogo de productos\n'
        '📋 Gestión de eventos escolares\n'
        '👥 Roles diferenciados: alumno / admin\n'
        '\n'
        'Transversal:\n'
        '🔒 Rutas protegidas — middleware JWT en cada endpoint sensible\n'
        '📱 Diseño responsivo — Mobile-first con Tailwind CSS\n'
        '🤖 IA integrada — Chatbot asistente en el dashboard',
        font_size=13, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Modelo de Sostenibilidad  (agregar modelo + costos lado izquierdo)
# ══════════════════════════════════════════════════════════════════════════════
s9 = slides[8]

# Modelo elegido
add_box(s9, 0.6, 2.2, 7.5, 0.55,
        'Modelo: B2B / Institución Educativa', font_size=17, bold=True, color=RED)
add_box(s9, 0.6, 2.85, 7.5, 0.9,
        'La escuela es el cliente. La plataforma se entrega sin costo y se sostiene con recursos existentes de la institución.',
        font_size=12, color=WHITE)

# Costos
add_box(s9, 0.6, 4.0, 7.5, 0.5,
        '📊 Análisis de Costos', font_size=15, bold=True, color=WHITE)
add_box(s9, 0.6, 4.6, 7.5, 1.8,
        '• Hosting (Replit plan gratuito):         $0 / mes\n'
        '• Base de datos (SQLite local):           $0 / mes\n'
        '• Email (Gmail institucional):              $0 / mes\n'
        '• Dominio .edu.mx:                              ~$200 MXN / año\n'
        '─────────────────────────────────────────\n'
        '  TOTAL:                                              ~$200 MXN / año',
        font_size=12, color=WHITE)

# Impacto técnico
add_box(s9, 0.6, 6.7, 7.5, 0.5,
        '💡 Impacto técnico del modelo', font_size=14, bold=True, color=WHITE)
add_box(s9, 0.6, 7.3, 7.5, 1.2,
        'Al no requerir ingresos propios, el sistema simula pagos en lugar de integrar pasarelas reales. Si la escuela escala, basta agregar OXXO Pay o Stripe al módulo Pago.jsx sin reescribir el resto.',
        font_size=12, color=(210,210,210))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Conclusión  (quitar info falsa, agregar GitHub + Demo + Preguntas)
# ══════════════════════════════════════════════════════════════════════════════
s11 = slides[10]

# Reemplazar +123-456-7890
tb_phone = find_shape(s11, 'TextBox 18')
if tb_phone:
    replace_text(tb_phone, '🚀 DEMO EN VIVO', font_size=16, bold=True, color=RED)

# Reemplazar www.presentacion.com
tb_web = find_shape(s11, 'TextBox 19')
if tb_web:
    replace_text(tb_web,
        'Login → Dashboard → Tienda → Pago → Recibo PDF',
        font_size=12, color=WHITE)

# Reemplazar dirección física
tb_addr = find_shape(s11, 'TextBox 20')
if tb_addr:
    replace_text(tb_addr,
        '📂 github.com/[su-usuario]/cbtis258-financieros\n'
        '\n'
        '¿Preguntas?',
        font_size=13, bold=False, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# Guardar
# ══════════════════════════════════════════════════════════════════════════════
prs.save(path_out)
print(f"✅ Guardado en: {path_out}")
